import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import re
import matplotlib.pyplot as plt
import matplotlib.patches as patches

# --- 1. 데이터 파싱 및 트리 구조화 ---
def parse_line(text):
    text = str(text).strip()
    match = re.match(r'^([\d\.]+)', text)
    if match:
        code = match.group(1).rstrip('.')
        level = code.count('.') + 1
        return {'id_code': code, 'text': text, 'level': level}
    return None

def build_tree(data):
    nodes = {}
    root_nodes = []
    for item in data:
        code = item['id_code']
        node = {'code': code, 'text': item['text'], 'level': item['level'], 'children': []}
        nodes[code] = node
        parts = code.split('.')
        if len(parts) > 1:
            parent_code = ".".join(parts[:-1])
            if parent_code in nodes:
                nodes[parent_code]['children'].append(node)
            else:
                if item['level'] == 1: root_nodes.append(node)
        else:
            root_nodes.append(node)
    return root_nodes

# --- 2. 좌표 계산 로직 (그룹 간격 최적화) ---
def calculate_layout(root_nodes, config):
    layout_data = []
    wbs_w = config['wbs_w']
    wbs_h = config['wbs_h']
    l1_gap_x = config['l1_gap_x']
    l2_gap_x = config['l2_gap_x']
    
    # 레벨별 그룹 사이 간격 (황금비율 적용값)
    group_v_gaps = {
        2: config['v_gap_1_2'],
        3: config['v_gap_2_3'],
        4: config['v_gap_3_4'],
        'deep': config['v_gap_deep']
    }
    
    # 직계(부모-자식) 사이 아주 좁은 고정 간격
    tight_gap = 0.05 

    start_x = (33.8 - wbs_w) / 2
    start_y = (19.05 - wbs_h) / 2

    l1_count = len(root_nodes)
    if l1_count == 0: return []
    l1_width = (wbs_w - (l1_gap_x * (l1_count - 1))) / l1_count

    for i, l1 in enumerate(root_nodes):
        x_l1 = start_x + (i * (l1_width + l1_gap_x))
        y_l1 = start_y
        h_l1 = 1.2
        layout_data.append({'node': l1, 'x': x_l1, 'y': y_l1, 'w': l1_width, 'h': h_l1, 'level': 1})

        if l1['children']:
            l2_count = len(l1['children'])
            l2_width = (l1_width - (l2_gap_x * (l2_count - 1))) / l2_count
            
            # 1레벨 -> 2레벨은 첫 시작이므로 그룹 간격 적용
            current_y_for_l2 = y_l1 + h_l1 + group_v_gaps[2]

            for j, l2 in enumerate(l1['children']):
                # 형제 L2 사이에는 간격을 크게 주지 않고 (이미 X축으로 나뉘어 있으므로) 
                # 하지만 세로형 WBS라면 여기서 y를 벌려야 함. 현재는 2레벨까지 가로 전개.
                x_l2 = x_l1 + (j * (l2_width + l2_gap_x))
                y_l2 = current_y_for_l2
                h_l2 = 1.0
                layout_data.append({'node': l2, 'x': x_l2, 'y': y_l2, 'w': l2_width, 'h': h_l2, 'level': 2})

                # 3레벨 이하 재귀적 배치 함수
                def draw_descendants(parent_node, px, py, pw, ph, level):
                    nonlocal layout_data
                    last_y = py + ph
                    
                    for idx, child in enumerate(parent_node['children']):
                        # 부모와 '첫 번째 자식' 사이는 촘촘하게(tight_gap)
                        # '두 번째 형제' 부터는 그룹 간격(group_v_gap) 적용
                        if idx == 0:
                            current_gap = tight_gap
                        else:
                            # 레벨에 맞는 그룹 간격 선택
                            current_gap = group_v_gaps.get(level + 1, group_v_gaps['deep'])
                        
                        target_y = last_y + current_gap
                        
                        # 너비 축소 및 우측 정렬
                        reduction = 0.3 * (child['level'] - 2)
                        c_w = max(l2_width - reduction, 2.0)
                        c_x = (px + pw) - c_w
                        c_h = 0.8
                        
                        layout_data.append({
                            'node': child, 'x': c_x, 'y': target_y, 'w': c_w, 'h': c_h, 'level': child['level']
                        })
                        
                        # 자식의 자식들을 그리기 위해 재귀 호출 (여기서 반환된 y값이 이 그룹의 진짜 끝)
                        last_y = draw_descendants(child, c_x, target_y, c_w, c_h, child['level'])
                    
                    return last_y

                # 2레벨 아래로 3레벨부터 시작
                draw_descendants(l2, x_l2, y_l2, l2_width, h_l2, 2)
                    
    return layout_data

# --- 3. 미리보기 & 4. PPT 생성 (이전과 거의 동일, 레이아웃 데이터만 활용) ---
def draw_preview(layout_data):
    fig, ax = plt.subplots(figsize=(10, 5.6))
    ax.set_xlim(0, 33.8)
    ax.set_ylim(0, 19.05)
    ax.invert_yaxis()
    ax.add_patch(patches.Rectangle((0, 0), 33.8, 19.05, linewidth=1, edgecolor='black', facecolor='#f9f9f9', alpha=0.5))
    for item in layout_data:
        lvl = item['level']
        color = '#1f497d' if lvl == 1 else '#365f91' if lvl == 2 else '#d9d9d9'
        rect = patches.Rectangle((item['x'], item['y']), item['w'], item['h'], linewidth=1, edgecolor='white', facecolor=color)
        ax.add_patch(rect)
        display_text = item['node']['text'][:12]
        txt_color = 'white' if lvl <= 2 else 'black'
        ax.text(item['x'] + item['w']/2, item['y'] + item['h']/2, display_text, color=txt_color, fontsize=6, ha='center', va='center')
    ax.set_axis_off()
    st.pyplot(fig)

def generate_ppt(layout_data):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(33.8), Cm(19.05)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for item in layout_data:
        shp = slide.shapes.add_shape(1, Cm(item['x']), Cm(item['y']), Cm(item['w']), Cm(item['h']))
        lvl = item['level']
        shp.fill.solid()
        if lvl == 1:
            shp.fill.fore_color.rgb = RGBColor(31, 73, 125)
            f_size, f_bold, f_color, align = Pt(12), True, RGBColor(255, 255, 255), PP_ALIGN.CENTER
        elif lvl == 2:
            shp.fill.fore_color.rgb = RGBColor(54, 95, 145)
            f_size, f_bold, f_color, align = Pt(10), False, RGBColor(255, 255, 255), PP_ALIGN.CENTER
        else:
            c = min(210 + (lvl * 8), 250)
            shp.fill.fore_color.rgb = RGBColor(c, c, c+5)
            shp.line.color.rgb = RGBColor(180, 180, 180)
            f_size, f_bold, f_color, align = Pt(9), False, RGBColor(0, 0, 0), PP_ALIGN.LEFT
        tf = shp.text_frame
        tf.text = item['node']['text']
        p = tf.paragraphs[0]
        p.font.size, p.font.bold, p.font.color.rgb, p.alignment = f_size, f_bold, f_color, align
    return prs

# --- 5. Streamlit UI ---
st.set_page_config(page_title="WBS Grouping Designer", layout="wide")
st.sidebar.title("🎨 WBS 그룹 디자인")

with st.sidebar.expander("📏 전체 영역 (cm)", expanded=True):
    wbs_w = st.number_input("가로 너비", 10.0, 32.0, 30.0, 0.5)
    wbs_h = st.number_input("세로 높이", 5.0, 18.0, 15.0, 0.5)

with st.sidebar.expander("↔️ 좌우 간격 (cm)", expanded=True):
    l1_gap_x = st.number_input("대그룹(L1) 간격", 0.0, 10.0, 1.5, 0.1)
    l2_gap_x = st.number_input("소그룹(L2) 간격", 0.0, 5.0, 0.5, 0.1)

with st.sidebar.expander("↕️ 그룹 간 간격 (황금비율)", expanded=True):
    st.info("부모-첫 자식은 밀착되고, 형제 그룹 사이만 벌어집니다.")
    auto_golden = st.checkbox("황금비율 모드 사용", value=True)
    base_v_gap = st.number_input("기준 그룹 간격", 0.1, 5.0, 0.8, 0.1)

    if auto_golden:
        v_gap_1_2 = base_v_gap
        v_gap_2_3 = round(v_gap_1_2 * 0.618, 2)
        v_gap_3_4 = round(v_gap_2_3 * 0.618, 2)
        v_gap_deep = round(v_gap_3_4 * 0.618, 2)
    else:
        v_gap_1_2 = st.number_input("1→2 간격", 0.0, 5.0, 0.6)
        v_gap_2_3 = st.number_input("2→3 그룹 간격", 0.0, 5.0, 0.4)
        v_gap_3_4 = st.number_input("3→4 그룹 간격", 0.0, 5.0, 0.2)
        v_gap_deep = st.number_input("깊은 레벨 그룹 간격", 0.0, 5.0, 0.1)

config = {
    'wbs_w': wbs_w, 'wbs_h': wbs_h, 'l1_gap_x': l1_gap_x, 'l2_gap_x': l2_gap_x,
    'v_gap_1_2': v_gap_1_2, 'v_gap_2_3': v_gap_2_3, 'v_gap_3_4': v_gap_3_4, 'v_gap_deep': v_gap_deep
}

st.title("📊 WBS 그룹 정렬 디자이너")
uploaded_file = st.file_uploader("엑셀 또는 PPT 파일 업로드", type=["xlsx", "pptx"])

if uploaded_file:
    raw_data = []
    if uploaded_file.name.endswith("xlsx"):
        df = pd.read_excel(uploaded_file)
        for val in df.iloc[:, 0]:
            p = parse_line(val)
            if p: raw_data.append(p)
    else:
        input_prs = Presentation(uploaded_file)
        for s in input_prs.slides:
            for shp in s.shapes:
                if hasattr(shp, "text"):
                    p = parse_line(shp.text)
                    if p: raw_data.append(p)

    if raw_data:
        raw_data.sort(key=lambda x: [int(i) for i in x['id_code'].split('.')])
        tree = build_tree(raw_data)
        layout_data = calculate_layout(tree, config)
        st.subheader("🖼️ 실시간 디자인 미리보기")
        draw_preview(layout_data)
        if st.button("🚀 PPT 생성 및 다운로드", use_container_width=True):
            final_ppt = generate_ppt(layout_data)
            ppt_io = io.BytesIO()
            final_ppt.save(ppt_io)
            ppt_io.seek(0)
            st.download_button("🎁 PPT 파일 다운로드", ppt_io, "WBS_Grouped.pptx")
